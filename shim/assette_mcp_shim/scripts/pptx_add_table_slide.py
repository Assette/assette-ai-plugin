#!/usr/bin/env python3
"""Surgically append a data-object-bound TABLE slide to an EXISTING Smart Page pair.

This fills the gap the fabricator skill documents in Workflow D: its
``fabricate`` command is a FULL REGENERATE from a spec, which would destroy a
real, plugin-authored deck. This script instead edits a pulled
``(.pptx, .metadata.json)`` pair *in place* — it adds ONE new slide carrying ONE
table shell, and leaves every existing slide, shape, chart, image, theme and tag
untouched.

It emits the REAL Assette plugin tag vocabulary (SHELLNAME / OBJECTTAG=
SelfFabricatedBasicTable / COLUMNSMAPPING / ROWSMAPPING / ROWMAPPER /
ADVANCESETTINGS / DISPLAYSETTINGS / COLUMNSTYLEMAPPING / …) rather than the
fabricator's consolidated ``ShellDetails`` tag, because that is the shape the
production generation engine has actually been observed to consume. The default
ADVANCESETTINGS below is copied verbatim from a known-good plugin-authored
``columnrow`` basic table.

Column keys, data types, runtime parameters and design-time parameters are
auto-derived from ``DataObjects/<name>.json`` (the cache
``get_data_object_metadata`` writes) when present; every derived value can be
overridden by a flag.

It reuses the skill's battle-tested ``assette_smartpage.tag_parts.inject_tags``
for the OOXML tag wiring (idempotent, additive), and re-opens the result to
self-verify before returning.

Usage:
    python -m assette_mcp_shim.scripts.pptx_add_table_slide \
        --pptx "SmartPages/ASC factsheet New Copy/ASC factsheet New Copy.pptx" \
        --metadata "SmartPages/ASC factsheet New Copy/ASC factsheet New Copy.metadata.json" \
        --data-object HoldingDetails

    # explicit columns / rows / category / shell name:
    python -m assette_mcp_shim.scripts.pptx_add_table_slide --pptx <f> --metadata <f> \
        --data-object HoldingDetails \
        --columns ISSUEDISPLAYNAME,ISSUENAME,RISKCOUNTRY,ISSUECOUNTRY,PERCENTAGE \
        --rows headings-data --data-category Holdings --shell-name "HoldingDetails Table"

    --dry-run prints the planned shell + tags and writes nothing.

After it returns 0, push the pair back with mcp__assette__update_smart_page
(pptxPath + metadataPath, bump slideCount, extend the `layouts` property by one
`ppLayoutBlank`).

Exit 0 on success, 2 on usage / validation error.
"""

from __future__ import annotations

import argparse
import json
import re
import sys
import uuid
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu

# Make the fabricator's Python package importable so we reuse its tag
# injector verbatim. The script lives at
# ``<plugin>/shim/assette_mcp_shim/scripts/`` and the fabricator skill
# lives at ``<plugin>/skills/assette-pptx-authoring/`` — both
# anchored on ``runtime.locations.plugin_root()``.
from assette_mcp_shim.runtime.locations import plugin_root  # noqa: E402

_SKILL_DIR = plugin_root() / "skills" / "assette-pptx-authoring"
if not _SKILL_DIR.is_dir():
    raise RuntimeError(
        f"Could not locate the assette-pptx-authoring skill at "
        f"{_SKILL_DIR}. If the plugin layout changed, update "
        "runtime.locations.plugin_root()."
    )
sys.path.insert(0, str(_SKILL_DIR))
from assette_smartpage.tag_parts import TagSpec, inject_tags  # noqa: E402

P = "http://schemas.openxmlformats.org/presentationml/2006/main"
R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

NUMERIC_TYPES = {"decimal", "number", "int", "integer", "currency", "percentage", "double", "float"}

# Known-good plugin default for a columnrow SelfFabricatedBasicTable.
DEFAULT_ADVANCE_SETTINGS = {
    "MappingType": "columnrow", "RowOrder": "receivedorder", "ColOrder": None,
    "MappingRowType": None, "Layout": "nonpaginated", "SiblingCount": 3,
    "IsHeadingShowEachFirstLevel": True, "IsHeadingShowEachCategory": True,
    "IsStartEachCategoryInFirstLevel": False, "IsChangedNewSlideHeight": None,
    "IsSwitchHeadingAndCategory": None, "IsChangedSecoundSlideHeight": False,
    "IsShowTitleInSubSequentPages": False, "HideRowsFromSecondTable": None,
    "SecondPageX": 0.0, "SecondPageY": 0.0, "SecondTableX": None, "SecondTableY": None,
    "GapBetweenLeftRight": 0.0, "GapBetweenTopBottom": 0.0, "Direction": "toptobottom",
    "MaintainFixedWidth": False, "FixedWidthColumns": None, "HiddenRowTypes": None,
    "IsShowHeadingInSubSequentPages": True, "FirstPageTitleRows": [],
    "HideColumnsFromSecondTable": [], "SubSequentPageTitleTrasnformType": "prefix",
    "SubSequentPagesTitleRows": [], "SubSequentPageTitleTrasnformText": "",
    "SubSequentPageTitleTrasnformRow": "none", "HideColumnAfterLevel": False,
    "IsBandedColor": False, "AlternativeColor1": "000000", "AlternativeColor2": "FFFFFF",
    "IsBandedColumnColor": False, "AlternativeColor3": "000000", "AlternativeColor4": "FFFFFF",
    "TocLevel0Format": None, "TocLevel1Format": None, "TocLevel2Format": None,
    "IsOverrideMargin": False, "MarginTop": None, "MarginBottom": None,
    "MarginLeft": None, "MarginRight": None, "ApplyBrandTheme": False,
    "KeepRowHeight": False, "KeepColumnWidth": False,
}

ROW_PRESETS = {
    "headings-data": [
        ("Headings", "Heading", "Headings", False),
        ("Data", "Data", "Data", False),
    ],
    "headings-data-summary": [
        ("Headings", "Heading", "Headings", False),
        ("Data", "Data", "Data", False),
        # Summary is an aggregate row the engine BINDS from the data object's
        # own Summary row (specialText "Total" + per-column totals); it must be
        # IsStatic=False or the engine treats it as static template content and
        # leaves it blank. (Matches a real plugin-authored deck, e.g. "MCP Test 21".)
        ("Summary", "Summary", "Summary", False),
    ],
}


def find_data_objects_dir(start: Path) -> Path | None:
    cur = start.resolve()
    for d in [cur, *cur.parents]:
        cand = d / "DataObjects"
        if cand.is_dir():
            return cand
    return None


def load_data_object(name: str, start: Path):
    do_dir = find_data_objects_dir(start)
    if not do_dir:
        return None
    f = do_dir / f"{name}.json"
    if not f.is_file():
        return None
    return json.loads(f.read_text(encoding="utf-8"))


def derive_columns(dom, override_keys):
    """Return (keys, headers, types) ordered."""
    cols = sorted(dom.get("Columns", []), key=lambda c: c.get("Order", 0)) if dom else []
    by_key = {(c.get("Key") or c.get("Name")): c for c in cols}
    if override_keys:
        keys = override_keys
    else:
        keys = [(c.get("Key") or c.get("Name")) for c in cols if (c.get("Key") or c.get("Name"))]
    headers, types = [], []
    for k in keys:
        c = by_key.get(k, {})
        headers.append(c.get("Name") or k)
        dt = (c.get("DataType") or "string").lower()
        types.append("decimal" if dt in NUMERIC_TYPES else "string")
    return keys, headers, types


def derive_params(dom):
    """Return (runtime_param_names, design_time_params_dict)."""
    runtime, design = [], {}
    for ip in (dom.get("InputParameters", []) if dom else []):
        nm = ip.get("Name")
        if not nm:
            continue
        if ip.get("RuntimeParameter"):
            runtime.append(nm)
        else:
            dv = ip.get("DefaultValue")
            design[nm] = "" if dv is None else str(dv)
    return runtime, design


def existing_instance_id(pptx: Path):
    """Reuse the INSTANCEID shared by existing slides if present."""
    zf = zipfile.ZipFile(str(pptx))
    try:
        for n in zf.namelist():
            if n.startswith("ppt/tags/") and n.endswith(".xml"):
                root = etree.fromstring(zf.read(n))
                for t in root.findall(f"{{{P}}}tag"):
                    if t.get("name") == "INSTANCEID" and t.get("val"):
                        return t.get("val")
    finally:
        zf.close()
    return None


def design_layout_theme(meta):
    """Pull DesignName / LayoutName from the first fabricated slide; ThemeId from top."""
    design, layout = "", "Blank"
    for s in meta.get("Slides", []):
        if s.get("IsFebricated") and s.get("DesignName"):
            design, layout = s.get("DesignName", ""), s.get("LayoutName", "Blank")
            break
    return design, layout, meta.get("ThemeId", "")


# lxml imported lazily so the module also loads where only python-pptx is needed.
from lxml import etree  # noqa: E402


def main(argv=None) -> int:
    ap = argparse.ArgumentParser(description="Append a data-object-bound table slide to an existing Smart Page pair.")
    ap.add_argument("--pptx", required=True)
    ap.add_argument("--metadata", required=True)
    ap.add_argument("--data-object", required=True, help="Data object name to bind the table to.")
    ap.add_argument("--columns", help="Comma-separated column KEYS (default: all visible columns from DataObjects/<name>.json).")
    ap.add_argument("--headers", help="Comma-separated header labels (default: column display names).")
    ap.add_argument("--rows", choices=sorted(ROW_PRESETS), default="headings-data")
    ap.add_argument("--shell-name", help="Shell display name (default: '<DataObject> Table').")
    ap.add_argument("--data-category", help="DATACATEGORY / BrandTheme.DataCategory (default: the data object name).")
    ap.add_argument("--runtime-params", help="Comma-separated runtime params (default: from data object metadata).")
    ap.add_argument("--design-name", help="Override slide DesignName (default: from an existing fabricated slide).")
    ap.add_argument("--layout-name", help="Override slide LayoutName (default: from an existing fabricated slide).")
    ap.add_argument("--dry-run", action="store_true", help="Print the planned shell + tags and write nothing.")
    args = ap.parse_args(argv)

    pptx = Path(args.pptx)
    meta_path = Path(args.metadata)
    if not pptx.is_file():
        print(f"error: pptx not found: {pptx}", file=sys.stderr); return 2
    if not meta_path.is_file():
        print(f"error: metadata not found: {meta_path}", file=sys.stderr); return 2

    dom = load_data_object(args.data_object, pptx.parent)
    override_keys = [c.strip() for c in args.columns.split(",")] if args.columns else None
    keys, headers, types = derive_columns(dom, override_keys)
    if args.headers:
        headers = [h.strip() for h in args.headers.split(",")]
    if not keys:
        print(f"error: no columns for '{args.data_object}'. Pass --columns or ensure "
              f"DataObjects/{args.data_object}.json exists.", file=sys.stderr)
        return 2
    if len(headers) != len(keys):
        print("error: --headers count must match column count.", file=sys.stderr); return 2

    do_runtime, design_params = derive_params(dom)
    runtime = [p.strip() for p in args.runtime_params.split(",")] if args.runtime_params else do_runtime
    shell_name = args.shell_name or f"{args.data_object} Table"
    data_category = args.data_category or args.data_object

    meta = json.loads(meta_path.read_text(encoding="utf-8"))
    # A slide-less base (e.g. fabricated with `slides: []` to inherit a Brand
    # Theme's masters) omits the "Slides" key entirely, which would crash the
    # append below AFTER the pptx was already mutated. Normalize it up-front.
    meta.setdefault("Slides", [])
    design, layout, _theme = design_layout_theme(meta)
    design = args.design_name or design
    layout = args.layout_name or layout

    # Build the mappings.
    ncols = len(keys)
    colmap = {str(i): keys[i] for i in range(ncols)}
    colstyle = [{"Tag": types[i], "ColumnIndex": i, "IsStatic": False, "IsVisible": False, "DisplayIndex": 0}
                for i in range(ncols)]
    rowsmap, rowmapper = [], {}
    for idx, (tag, cat, style, is_static) in enumerate(ROW_PRESETS[args.rows]):
        rowsmap.append({"Tag": tag, "RowIndex": idx, "RowCategory": cat, "IsStatic": is_static,
                        "Value": tag, "DisplayIndex": 0, "RowStyle": style})
        rowmapper[str(idx)] = tag
    display = {"ObjectName": args.data_object, "Type": "data", "Parameters": design_params,
               "CustomSettings": {"GeneralSettings": None, "ColumnSettings": {}, "SeriesSettings": None, "ImageSettings": None}}

    shell_id = str(uuid.uuid4())
    slide_guid = str(uuid.uuid4())

    shell = {
        "ShellId": shell_id, "Name": shell_name, "Source": None, "Type": "Table",
        "TableSettings": {
            "Type": "SelfFabricatedBasicTable", "Settings": display,
            "AdvanceSettings": DEFAULT_ADVANCE_SETTINGS,
            "RowsMapping": rowsmap, "ColumnsMapping": colmap,
            "GenerationSettings": [], "RuntimeParameters": runtime,
            "GenerationTimeParameters": [], "LockSettings": [], "CellsMapping": None,
            "IsTOCTable": False, "ColumnStyleMapping": colstyle,
        },
        "ChartSettings": None, "LabelSettings": None, "PictureSettings": None,
        "DisclosureSettings": None, "TOC": None, "PageMarker": None,
        "PageMarkerReferences": None, "TitleId": "", "FootnoteId": "",
        "BrandTheme": {"DataCategory": data_category, "DataObject": args.data_object,
                       "DataCategoryId": None, "ChartTypes": None, "IsExtendChart": False,
                       "ThemeVersion": None, "BackupStyles": None, "KeepRowHeight": False,
                       "KeepColumnWidth": False},
    }

    shape_tags = {
        "SHELLNAME": shell_name, "GENERATIONSETTINGS": "[]", "DISCLOSURESSCOPE": "{}",
        "LOCKSETTINGS": "[]", "ISTOCTABLE": "False", "SHAPETYPE": "Table",
        "UNIQUEID": shell_id, "TABLETYPENEW": "columnrow", "UNIQUEFORTHEMEID": shell_id,
        "OBJECTTAG": "SelfFabricatedBasicTable", "RUNTIMEPARAMETERS": json.dumps(runtime),
        "GENERATIONTIMESETTINGS": "[]", "ADVANCESETTINGS": json.dumps(DEFAULT_ADVANCE_SETTINGS),
        "DATACATEGORY": data_category, "COLUMNSMAPPING": json.dumps(colmap),
        "COLUMNSTYLEMAPPING": json.dumps(colstyle), "DATAOBJECT": args.data_object,
        "DISPLAYSETTINGS": json.dumps(display), "ROWMAPPER": json.dumps(rowmapper),
        "ROWSMAPPING": json.dumps(rowsmap),
    }

    print(f"data object : {args.data_object}  ({'metadata found' if dom else 'NO DataObjects/*.json — using flags only'})")
    print(f"columns     : {keys}")
    print(f"types       : {types}")
    print(f"rows        : {args.rows}  ({len(rowsmap)} rows)")
    print(f"runtime     : {runtime}")
    print(f"shell name  : {shell_name}   data category: {data_category}")
    print(f"design/layout: {design} / {layout}")

    if args.dry_run:
        print("\n--- DRY RUN: shell that would be appended to metadata.json ---")
        print(json.dumps(shell, indent=1, ensure_ascii=False))
        print("\n--- shape tag keys ---")
        print(", ".join(shape_tags))
        print("\n(no files written)")
        return 0

    # ---- mutate the pptx: add slide + table ----
    prs = Presentation(str(pptx))
    before = set(zipfile.ZipFile(str(pptx)).namelist())
    blank = next((l for l in prs.slide_layouts if (l.name or "").strip().lower() == "blank"), None)
    if blank is None:
        blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank)
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    W, H = prs.slide_width, prs.slide_height
    nrows = len(rowsmap)
    gf = slide.shapes.add_table(nrows, ncols, Emu(int(W * 0.04)), Emu(int(H * 0.16)),
                                Emu(int(W * 0.92)), Emu(int(H * 0.5)))
    shape_name = f"tbl_{shell_id}"
    gf.name = shape_name
    tbl = gf.table
    for c in range(ncols):
        tbl.cell(0, c).text = headers[c]
    if nrows > 1:
        for c in range(ncols):
            tbl.cell(1, c).text = "<data>"
    # Right-align every numeric column (header + body) so the generated table
    # renders numbers right-aligned; the engine clones each template cell's
    # paragraph alignment per column. `types[c]` is "decimal" for numerics.
    for c in range(ncols):
        if types[c] != "decimal":
            continue
        for r in range(nrows):
            for para in tbl.cell(r, c).text_frame.paragraphs:
                para.alignment = PP_ALIGN.RIGHT
    new_index = len(prs.slides)
    prs.save(str(pptx))

    after = set(zipfile.ZipFile(str(pptx)).namelist())
    new_parts = [p for p in (after - before) if re.match(r"ppt/slides/slide\d+\.xml$", p)]
    if len(new_parts) != 1:
        print(f"error: expected exactly one new slide part, got {new_parts}", file=sys.stderr); return 2
    slide_part = new_parts[0]

    instance_id = existing_instance_id(pptx) or str(uuid.uuid4())
    slide_tags = {"INSTANCEID": instance_id, "ASSETTESLIDEGUID": slide_guid, "ISFABRICATEDSLIDE": "true"}

    inject_tags(pptx, [
        TagSpec(target_path=slide_part,
                owner_xpath=(f".//p:sp[.//p:cNvPr[@name='{shape_name}']]/p:nvSpPr/p:nvPr"
                             f" | .//p:graphicFrame[.//p:cNvPr[@name='{shape_name}']]/p:nvGraphicFramePr/p:nvPr"),
                tags=shape_tags),
        TagSpec(target_path=slide_part, owner_xpath="//p:cSld", tags=slide_tags),
    ])

    # ---- update metadata.json ----
    new_slide = {
        "Sequence": new_index, "SlideId": slide_guid, "DesignName": design,
        "LayoutName": layout, "IsFebricated": True, "IsTocSlide": False,
        "Shells": [shell], "HasDynamicParameters": False,
        "SlideMargin": {"isOverride": False, "Top": 0.0, "Bottom": 0.0, "Left": 0.0, "Right": 0.0},
        "AdjustObjects": [],
    }
    # Defense-in-depth: a slide-less base (spec `slides: []`) historically
    # emitted a metadata.json with NO "Slides" key (the serializer's clean pass
    # dropped the empty list). The root fix now always emits "Slides": [], but
    # keep this guard so an older/hand-written pair can't KeyError here.
    meta.setdefault("Slides", [])
    meta["Slides"].append(new_slide)
    meta["NoOfSlides"] = len(meta["Slides"])
    rp = meta.setdefault("RuntimeParameters", [])
    for p in runtime:
        if p not in rp:
            rp.append(p)
    meta_path.write_text(json.dumps(meta, indent=1, ensure_ascii=False), encoding="utf-8")

    # ---- self-verify ----
    prs2 = Presentation(str(pptx))
    assert len(prs2.slides) == new_index, "slide count mismatch after save"
    zf = zipfile.ZipFile(str(pptx))
    rels = {}
    d = Path(slide_part)
    for rel in etree.fromstring(zf.read(f"{d.parent.as_posix()}/_rels/{d.name}.rels")):
        if (rel.get("Type") or "").endswith("/tags"):
            rels[rel.get("Id")] = rel.get("Target").replace("../", "ppt/")
    root = etree.fromstring(zf.read(slide_part))
    found_uid = False
    for gf2 in root.iter(f"{{{P}}}graphicFrame"):
        nv = gf2.find(f".//{{{P}}}nvPr")
        cdl = nv.find(f"{{{P}}}custDataLst") if nv is not None else None
        if cdl is None:
            continue
        for t in cdl.findall(f"{{{P}}}tags"):
            part = rels.get(t.get(f"{{{R}}}id"))
            if part:
                vals = {x.get("name"): x.get("val") for x in etree.fromstring(zf.read(part)).findall(f"{{{P}}}tag")}
                if vals.get("UNIQUEID") == shell_id:
                    found_uid = True
    zf.close()
    assert found_uid, "new table UNIQUEID tag not found after injection"

    print(f"\nOK: added slide {new_index} ({slide_part}); shell_id={shell_id}")
    print(f"    pptx     : {pptx}")
    print(f"    metadata : {meta_path}  (NoOfSlides={meta['NoOfSlides']})")
    print(f"    NEXT     : update_smart_page id=<id> pptxPath/metadataPath, slideCount={new_index}, "
          f"extend `layouts` by one ppLayoutBlank.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
